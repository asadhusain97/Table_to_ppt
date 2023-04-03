from flask import render_template, request, redirect, url_for, send_file
import os
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
import io
import requests

from app import app


def transform_file(file_path):
    # Load the question file as pd.DataFrame
    question_df = pd.read_csv(file_path, sep="\t")
    max_rows = len(question_df)
    max_cols = len(question_df.columns)

    # Path where you want to save your new PowerPoint presentation
    # assuming file name includes something similar to ZQL_S3_GW7 Question Set.tsv
    ppt_path = file_path[:-27]  # same path/ folder as the question file
    ppt_name = file_path[-27:-17] + " Deck.pptx"  # Deck in place of Question Set
    powerpoint_path = ppt_path + ppt_name

    # Define the position and size of each cell on the slide
    cell_width = Inches(2)
    cell_height = Inches(1)
    cell_padding = Inches(0.1)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    slide_margin_left = Inches(0.5)
    slide_margin_top = Inches(0.5)
    max_cell_width = slide_width - 2 * slide_margin_left

    def add_text_to_slide(
            slide_var,
            cell_value_var,
            cell_left=slide_margin_left,
            cell_top=slide_margin_top,
            cell_width_var=cell_width,
            cell_height_var=cell_height,
            font_size=24,
            font_bold=True,
            word_wrap=True,
    ):
        textbox = slide_var.shapes.add_textbox(cell_left, cell_top, cell_width_var, cell_height_var)
        paragraph = textbox.text_frame.add_paragraph()
        paragraph.text = str(cell_value_var)
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = font_bold
        textbox.text_frame.word_wrap = word_wrap

    def add_image_to_slide(
            slide_var,
            cell_value_var,
            cell_left=slide_margin_left,
            cell_top=slide_margin_top,
            image_width=5,
            image_height=5,
    ):
        try:
            response = requests.get(cell_value_var)
        except ConnectionError:
            return 0
        image_data = response.content
        image_stream = io.BytesIO(image_data)
        picture = slide_var.shapes.add_picture(
            image_stream,
            left=cell_left,
            top=cell_top,
            width=Inches(image_width),
            height=Inches(image_height),
        )

    # Create a new PowerPoint presentation and add a slide for each row of data
    presentation = Presentation()
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    # Title page
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_text_to_slide(
        slide,
        "Zephyr Quiz",
        cell_left=slide_margin_left,
        cell_top=slide_margin_top + Inches(2),
        cell_width=max_cell_width,
        cell_height=cell_height,
        font_size=48,
        font_bold=True,
        word_wrap=False,
    )
    add_text_to_slide(
        slide,
        f"Season {ppt_name[5]} Game {ppt_name[9]}",
        cell_left=slide_margin_left,
        cell_top=slide_margin_top + Inches(4),
        cell_width=max_cell_width,
        cell_height=cell_height,
        font_size=30,
        font_bold=False,
        word_wrap=False,
    )
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for row in range(max_rows):
        for col in range(max_cols):
            cell_value = question_df.iloc[row, col]
            if col == 0:  # roundNo - ex. Round1 Player 1
                add_text_to_slide(
                    slide,
                    cell_value,
                    cell_left=slide_margin_left,
                    cell_top=slide_margin_top + Inches(0.25),
                    cell_width=max_cell_width,
                    cell_height=cell_height,
                    font_size=18,
                    font_bold=False,
                    word_wrap=False,
                )
            elif col == 1:  # questionNo - ex. Question1
                add_text_to_slide(
                    slide,
                    f"Question {row + 1}",
                    cell_left=slide_margin_left,
                    cell_top=slide_margin_top - Inches(0.25),
                    cell_width=max_cell_width,
                    cell_height=cell_height,
                    font_size=24,
                    font_bold=True,
                    word_wrap=False,
                )
            elif col == 2:  # questionText - ex. Who?'s the fairest of them all?
                add_text_to_slide(
                    slide,
                    cell_value,
                    cell_left=slide_margin_left,
                    cell_top=slide_margin_top + Inches(1.5),
                    cell_width=max_cell_width,
                    cell_height=cell_height,
                    font_size=24,
                    font_bold=True,
                    word_wrap=True,
                )
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                add_text_to_slide(
                    slide,
                    "SAFETY SLIDE",
                    cell_left=slide_margin_left,
                    cell_top=slide_margin_top - Inches(0.25),
                    cell_width=max_cell_width,
                    cell_height=cell_height,
                    font_size=24,
                    font_bold=True,
                    word_wrap=False,
                )
                safety_img_path = "https://i.imgur.com/I9kyepF.jpeg"  # "https://imgur.com/4InhUGd.gif"
                add_image_to_slide(
                    slide,
                    safety_img_path,
                    cell_left=slide_margin_left + Inches(2),
                    cell_top=slide_margin_top + Inches(1),
                    image_width=5,
                    image_height=5,
                )
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            elif col == 3:  # imageUrl - ex. linkto@picture.com
                add_text_to_slide(
                    slide,
                    "ANSWER SLIDE",
                    cell_left=slide_margin_left,
                    cell_top=slide_margin_top - Inches(0.25),
                    cell_width=max_cell_width,
                    cell_height=cell_height,
                    font_size=24,
                    font_bold=True,
                    word_wrap=False,
                )
                if cell_value == cell_value:  # nan - no link to image
                    add_image_to_slide(
                        slide,
                        cell_value,
                        cell_left=slide_width / 2,
                        cell_top=slide_margin_top + Inches(1.5),
                        image_width=4,
                        image_height=4,
                    )
            elif col == 4:  # answerText - ex. Snow White
                if "<br> <img src" in cell_value:
                    ans_txt = cell_value[: cell_value.find(" <br>")]
                    img_path = cell_value[
                               cell_value.find("https"): cell_value.find("https") + 31
                               ]  # length of link fixed - 31; png, jpg
                    add_text_to_slide(
                        slide,
                        ans_txt,
                        cell_left=slide_margin_left,
                        cell_top=slide_margin_top + Inches(1.5),
                        cell_width=max_cell_width / 2,
                        cell_height=cell_height,
                        font_size=24,
                        font_bold=True,
                        word_wrap=True,
                    )
                    add_image_to_slide(
                        slide,
                        img_path,
                        cell_left=slide_margin_left,
                        cell_top=slide_margin_top + Inches(3.5),
                        image_width=2.5,
                        image_height=2.5,
                    )
                else:
                    add_text_to_slide(
                        slide,
                        cell_value,
                        cell_left=slide_margin_left,
                        cell_top=slide_margin_top + Inches(1.5),
                        cell_width=max_cell_width / 2,
                        cell_height=cell_height,
                        font_size=24,
                        font_bold=True,
                        word_wrap=True,
                    )
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    # End page
    add_text_to_slide(
        slide,
        "Fin",
        cell_left=slide_margin_left,
        cell_top=slide_margin_top + Inches(2),
        cell_width=max_cell_width,
        cell_height=cell_height,
        font_size=48,
        font_bold=True,
        word_wrap=False,
    )

    # Save the PowerPoint presentation to the specified path
    presentation.save(powerpoint_path)
    return powerpoint_path


@app.route('/')
def home():
    return render_template('home.html')


@app.route('/upload', methods=['POST'])
def upload():
    if request.method == 'POST':
        file = request.files['file']
        file.save(os.path.join('uploads', file.filename))
        return render_template('home.html', success='Success!')
    else:
        return render_template('home.html')


@app.route('/download/<filename>')
def download(filename):
    try:
        print(filename)
        return send_file(os.path.join(app.upload_folder, filename), attachment_filename=filename[:-3], as_attachment=True)
    except Exception as e:
        return str(e)


#
# @app.route('/')
# def home():
#     return render_template('home.html')
#
#
# @app.route('/uploads', methods=['GET', 'POST'])
# def upload_file():
#     if request.method == 'POST':
#         file = request.files['file']
#         filename = file.filename
#         file.save(os.path.join(app.upload_folder, filename))
#         return redirect(url_for('uploaded_file',
#                                 filename=filename))
#     return render_template('upload.html')
#
#
# @app.route('/uploads/<filename>')
# def uploaded_file(filename):
#     return f'{filename} is uploaded!!!'
#
#
# @app.route('/uploads/<filename>')
# def download(filename):
#     # Get the path to the uploaded file
#     file_path = os.path.join(app.upload_folder, filename)
#
#     # Check if the file exists
#     if not os.path.exists(file_path):
#         return "File not found", 404
#
#     # Send the file to the user
#     return send_file(file_path, as_attachment=True)
